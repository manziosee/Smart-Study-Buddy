import os
import zipfile
from docx import Document
from pptx import Presentation
import json

def extract_text_from_docx(file_path):
    """Extract text from DOCX files"""
    try:
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text.strip())
        return '\n'.join(text)
    except Exception as e:
        raise Exception(f"Failed to extract text from DOCX: {str(e)}")

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint files"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())
        return '\n'.join(text)
    except Exception as e:
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")

def extract_text_from_json(file_path):
    """Extract text from JSON files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        def extract_strings(obj):
            if isinstance(obj, dict):
                return ' '.join([extract_strings(v) for v in obj.values()])
            elif isinstance(obj, list):
                return ' '.join([extract_strings(item) for item in obj])
            elif isinstance(obj, str):
                return obj
            else:
                return str(obj)
        
        return extract_strings(data)
    except Exception as e:
        raise Exception(f"Failed to extract text from JSON: {str(e)}")

def process_advanced_file(file):
    """Process advanced file formats"""
    file_extension = os.path.splitext(file.name)[1].lower()
    temp_path = f"/tmp/{file.name}"
    
    # Save file temporarily
    with open(temp_path, 'wb+') as temp_file:
        for chunk in file.chunks():
            temp_file.write(chunk)
    
    try:
        if file_extension == '.docx':
            text = extract_text_from_docx(temp_path)
        elif file_extension == '.pptx':
            text = extract_text_from_pptx(temp_path)
        elif file_extension == '.json':
            text = extract_text_from_json(temp_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        os.remove(temp_path)
        return text
        
    except Exception as e:
        if os.path.exists(temp_path):
            os.remove(temp_path)
        raise e

SUPPORTED_FORMATS = ['.pdf', '.txt', '.md', '.docx', '.pptx', '.json']